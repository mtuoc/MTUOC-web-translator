#pip install python-docx
#pip install python-pptx
#pip install odfpy

from docx import Document
from pptx import Presentation

class DocumentCleaner:
    def __init__(self):
        pass
        
    def cleanDOCX(self, input_file: str, output_file: str):
        try:
            doc = Document(input_file)            
            for paragraph in doc.paragraphs:
                cleaned_runs = []
                prev_run = None
                
                for run in paragraph.runs:
                    if prev_run and run.bold == prev_run.bold and run.italic == prev_run.italic \
                            and run.underline == prev_run.underline and run.font.size == prev_run.font.size:
                        prev_run.text += run.text
                    else:
                        cleaned_runs.append(run)
                        prev_run = run
                
                paragraph.clear()
                for cleaned_run in cleaned_runs:
                    new_run = paragraph.add_run(cleaned_run.text)
                    new_run.bold = cleaned_run.bold
                    new_run.italic = cleaned_run.italic
                    new_run.underline = cleaned_run.underline
                    new_run.font.size = cleaned_run.font.size

            doc.save(output_file)
        except Exception as e:
            print(f"An error occurred: {e}")

    def cleanPPTX(self, input_file: str, output_file: str):
        prs = Presentation(input_file)
        for slide in prs.slides:
            # Loop through all shapes in the slide
            for shape in slide.shapes:
                # Check if the shape has text (only TextBoxes, Titles, etc.)
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    # Loop through paragraphs in the text frame
                    for paragraph in shape.text_frame.paragraphs:
                        cleaned_runs = []
                        prev_run = None

                        # Loop through all the runs in the paragraph
                        for run in paragraph.runs:
                            # Compare formatting of consecutive runs
                            if prev_run and run.bold == prev_run.bold and run.italic == prev_run.italic \
                                    and run.underline == prev_run.underline and run.font.size == prev_run.font.size:
                                # If formatting is the same, append the text
                                prev_run.text += run.text
                            else:
                                # If formatting differs, create a new cleaned run
                                cleaned_runs.append(run)
                                prev_run = run

                        # Clear existing runs and re-add the cleaned runs
                        paragraph.clear()  # Removes existing runs
                        for cleaned_run in cleaned_runs:
                            new_run = paragraph.add_run(cleaned_run.text)
                            new_run.bold = cleaned_run.bold
                            new_run.italic = cleaned_run.italic
                            new_run.underline = cleaned_run.underline
                            new_run.font.size = cleaned_run.font.size

        # Save the cleaned presentation
        prs.save(output_file)
    